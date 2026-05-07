"""
Add four deeper-dive slides on the nonlinear optimizers and reorder them
into the right place in the deck.

Slides (inserted after the existing "Three Optimizers, One Interface"):
  - Gradient Descent — Steepest Descent on SSE
  - Gauss-Newton with Self-Activating LM Damping
  - BFGS — Quasi-Newton with Line Search
  - Optimizer Trajectories on a 2D Problem  (uses optimizer_trajectories.png)

Idempotent: re-runs are no-ops once the slides are present and ordered.

Usage:  python scripts/build_optimizer_slides.py
"""

from __future__ import annotations

import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT / "scripts"))

from build_slides import (        # noqa: E402  (path setup must come first)
    PPTX_PATH, RESULTS_DIR, CYAN,
    add_one_card_slide, add_two_card_slide, add_image_with_bullets_slide,
    _existing_titles,
)
from pptx import Presentation     # noqa: E402
from pptx.oxml.ns import qn       # noqa: E402


# ---------------------------------------------------------------------------
# Slide content
# ---------------------------------------------------------------------------

GD_TITLE = "Gradient Descent — Steepest Descent on SSE"
GN_TITLE = "Gauss-Newton with Self-Activating LM Damping"
BFGS_TITLE = "BFGS — Quasi-Newton with Line Search"
TRAJ_TITLE = "Optimizer Trajectories on a 2D Problem"

# Anchor: the existing overview slide that the new slides should sit AFTER.
ANCHOR_TITLE = "Three Optimizers, One Interface"


def add_gd_slide(prs):
    add_two_card_slide(
        prs,
        title=GD_TITLE,
        left_card_title="How it works",
        left_bullets=[
            "Loss:  L(θ) = (1/n) · Σᵢ (yᵢ − f(xᵢ, θ))²",
            "Gradient (with r = y − f(X, θ)):  ∇L(θ) = −(2/n) · Jᵀ r",
            "Update:  θₖ₊₁ = θₖ − αₖ · ∇L(θₖ)",
            "Step size:  αₖ = α₀ · decayᵏ  (optional shrinkage)",
            "Convergence:  |Lₖ − Lₖ₊₁| / Lₖ < tol",
            "Pure first-order method — only the gradient is consulted",
        ],
        right_card_title="When to use it (and when not to)",
        right_bullets=[
            "Strengths: simplest possible; works without Hessian; easy to reason about",
            "Weaknesses: only linear convergence — typically 10–100× more iterations than GN/BFGS",
            "Sensitive to learning rate: too large overshoots; too small crawls",
            "Struggles when parameters have very different scales (ill-conditioning)",
            "In our suite: pedagogical baseline only",
            "Code:  ta.GradientDescentRegressor(model_fn=f, learning_rate=1e-4, decay=0.9999, max_iter=5000)",
        ],
        left_accent=CYAN, right_accent=CYAN,
    )


def add_gn_slide(prs):
    add_two_card_slide(
        prs,
        title=GN_TITLE,
        left_card_title="How it works",
        left_bullets=[
            "Linearize the model around current θ:  f(X, θ + δ) ≈ f(X, θ) + J · δ",
            "Local SSE becomes quadratic in δ — solve a least-squares problem each step",
            "Damped normal equation:  (Jᵀ J + λ I) · δ = Jᵀ r",
            "Update:  θ ← θ + δ",
            "Damping limits — λ = 0 → pure Gauss-Newton;  λ → ∞ → ≈ gradient descent with step 1/λ",
            "Self-activating LM: λ starts at 0; on rejected step λ ← max(λ·factor, 1e-4); on accepted step λ ← λ/factor",
        ],
        right_card_title="When to use it (and when not to)",
        right_bullets=[
            "Strengths: superlinear / quadratic convergence near a good init",
            "LM auto-activation makes it robust without manual λ tuning",
            "Weaknesses: needs Jᵀ J to be (nearly) invertible",
            "Sensitive to initialization when far from the optimum — the linearization stops being accurate",
            "Best when residuals are small and the model is smooth",
            "Code:  ta.GaussNewtonRegressor(model_fn=f, max_iter=100, damping=0.0, damping_factor=10.0)",
        ],
        left_accent=CYAN, right_accent=CYAN,
    )


def add_bfgs_slide(prs):
    add_two_card_slide(
        prs,
        title=BFGS_TITLE,
        left_card_title="How it works",
        left_bullets=[
            "Quasi-Newton: maintain an inverse-Hessian approximation Hₖ; update it from successive gradients",
            "Search direction:  dₖ = − Hₖ · ∇L(θₖ)",
            "Step size αₖ via backtracking Armijo line search:  L(θ + αd) ≤ L(θ) + c₁ α (∇L)ᵀd",
            "Position update:  θₖ₊₁ = θₖ + αₖ dₖ",
            "Curvature update — with sₖ = αₖdₖ, yₖ = ∇Lₖ₊₁ − ∇Lₖ, ρₖ = 1 / (yₖᵀsₖ):",
            "    Hₖ₊₁ = (I − ρₖ sₖ yₖᵀ) · Hₖ · (I − ρₖ yₖ sₖᵀ) + ρₖ sₖ sₖᵀ",
        ],
        right_card_title="When to use it (and when not to)",
        right_bullets=[
            "Strengths: superlinear convergence WITHOUT computing the true Hessian",
            "Robust default — works on a wide range of well-formulated problems",
            "BFGS update is skipped when sᵀy ≤ ε (curvature condition violated)",
            "Symmetry enforced after each update: H ← (H + Hᵀ) / 2",
            "Reset H = I if ‖H‖ ever blows up — protects against drift",
            "Code:  ta.BFGSRegressor(model_fn=f, max_iter=500, c1=1e-4)",
        ],
        left_accent=CYAN, right_accent=CYAN,
    )


def add_trajectory_slide(prs):
    add_image_with_bullets_slide(
        prs,
        title=TRAJ_TITLE,
        image_path=RESULTS_DIR / "optimizer_trajectories.png",
        image_position="below",
        bullets_card_title="Same problem, same start, three trajectories",
        bullets=[
            "Problem: y = a · exp(b · x)  on exponential_multiplicative; θ₀ = (1, 0.1) (black X); truth (★) = (a=2, b=0.7)",
            "Gradient Descent — hundreds of tiny steps that snake along the SSE basin without reaching the optimum",
            "Gauss-Newton + LM — ~8 bold steps; near-quadratic convergence once inside the basin",
            "BFGS — ~17 steps; the line search adapts the step length each iteration → curved descent",
            "All three reach essentially the same RMSE; cost differs by 10–100×",
        ],
    )


# ---------------------------------------------------------------------------
# Reordering
# ---------------------------------------------------------------------------

def _slide_titles(prs):
    return _existing_titles(prs)


def reorder_to_anchor(prs, target_titles, anchor_title):
    """Move slides whose first-paragraph title is in `target_titles` so they
    appear directly after the slide whose title is `anchor_title`.  Order
    among the targets is preserved as given in `target_titles`."""
    titles = _slide_titles(prs)
    if anchor_title not in titles:
        print(f"warn: anchor {anchor_title!r} not found; skipping reorder")
        return False

    # We need to reorder via the underlying _sldIdLst element.
    sldIdLst = prs.slides._sldIdLst
    sld_ids = list(sldIdLst)

    # Map title -> sldId element (zero-indexed), using the first match.
    by_title = {}
    for sid, t in zip(sld_ids, titles):
        if t and t not in by_title:
            by_title[t] = sid

    anchor_sid = by_title[anchor_title]
    anchor_idx = sld_ids.index(anchor_sid)

    # Pull each target out of the list, preserving requested order.
    moving = []
    for t in target_titles:
        sid = by_title.get(t)
        if sid is None:
            print(f"warn: target {t!r} not found; skipping that one")
            continue
        moving.append(sid)

    # Check whether they are already in the desired order right after anchor.
    expected = [anchor_sid] + moving
    actual_window = sld_ids[anchor_idx : anchor_idx + len(expected)]
    if actual_window == expected:
        return False  # already in place

    for sid in moving:
        sldIdLst.remove(sid)
    # After removals, recompute anchor index.
    sld_ids = list(sldIdLst)
    anchor_idx = sld_ids.index(anchor_sid)
    # Insert in order, just after anchor.
    for offset, sid in enumerate(moving, start=1):
        sldIdLst.insert(anchor_idx + offset, sid)
    return True


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation(PPTX_PATH)
    titles_before = _slide_titles(prs)
    print(f"Opened {PPTX_PATH.name} ({len(prs.slides)} existing slides)")

    added = 0
    if GD_TITLE not in titles_before:
        add_gd_slide(prs);   added += 1
    if GN_TITLE not in titles_before:
        add_gn_slide(prs);   added += 1
    if BFGS_TITLE not in titles_before:
        add_bfgs_slide(prs); added += 1
    if TRAJ_TITLE not in titles_before:
        add_trajectory_slide(prs); added += 1
    print(f"Appended {added} new slide(s).")

    moved = reorder_to_anchor(
        prs,
        target_titles=[GD_TITLE, GN_TITLE, BFGS_TITLE, TRAJ_TITLE],
        anchor_title=ANCHOR_TITLE,
    )
    if moved:
        print(f"Reordered: 4 slides moved to sit right after '{ANCHOR_TITLE}'.")
    else:
        print("No reorder needed.")

    if added or moved:
        prs.save(PPTX_PATH)
        print(f"Saved: {PPTX_PATH}")
    else:
        print("Deck unchanged.")


if __name__ == "__main__":
    main()
