"""
Regenerate cost_pareto.png AND the five per-dataset Pareto PNGs from the
cached cost_summary.csv (no rerun of the benchmark).  Then replace the
picture on the existing "Cost vs. Accuracy" slide with the updated
combined image, preserving its position and size.

Per-dataset PNGs land in comparison_results/cost_pareto_<dataset>.png and
are NOT auto-inserted into the deck — copy them into slides as you like.

Usage:  python scripts/refresh_pareto_slide.py
"""

from __future__ import annotations

import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from cost_analysis import plot_pareto, plot_pareto_per_dataset, OUTPUT_DIR

PPTX_PATH = PROJECT_ROOT / "ZZU_Transformations_Presentation.pptx"
SUMMARY_CSV = OUTPUT_DIR / "cost_summary.csv"
PARETO_PNG = OUTPUT_DIR / "cost_pareto.png"
SLIDE_TITLE = "Cost vs. Accuracy"


def regenerate_png():
    if not SUMMARY_CSV.exists():
        sys.exit(f"error: {SUMMARY_CSV} not found — run cost_analysis.py first")
    summary = pd.read_csv(SUMMARY_CSV)
    plot_pareto(summary, PARETO_PNG)
    print(f"Regenerated {PARETO_PNG.name}")
    paths = plot_pareto_per_dataset(summary, OUTPUT_DIR)
    print("Per-dataset PNGs:")
    for p in paths:
        print(f"  - {p.name}")


def find_slide_by_title(prs, title):
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t.split("\n")[0].strip() == title:
                    return i, slide
    return None, None


def replace_picture_on_slide(slide, new_image_path):
    """Swap the (first) picture on the slide for `new_image_path`,
    preserving its left/top/width/height."""
    for sh in list(slide.shapes):
        if sh.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            left, top, width, height = sh.left, sh.top, sh.width, sh.height
            sp = sh._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(
                str(new_image_path), left, top, width=width, height=height
            )
            return True
    return False


def main():
    regenerate_png()

    prs = Presentation(PPTX_PATH)
    idx, slide = find_slide_by_title(prs, SLIDE_TITLE)
    if slide is None:
        sys.exit(f"error: slide titled {SLIDE_TITLE!r} not found")

    if replace_picture_on_slide(slide, PARETO_PNG):
        prs.save(PPTX_PATH)
        print(f"Replaced picture on slide {idx + 1} ('{SLIDE_TITLE}'); saved deck.")
    else:
        print(f"warn: no picture shape found on slide {idx + 1} — nothing replaced.")


if __name__ == "__main__":
    main()
