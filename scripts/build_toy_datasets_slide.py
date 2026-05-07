"""
Build a standalone .pptx with one slide introducing the five synthetic
toy datasets used in the ZZU comparison.  The slide is in a separate
file so it can be copied into the online presentation manually.

Outputs (project root):
  - toy_datasets_panel.png        2x3 panel of the five datasets
  - toy_datasets_intro.pptx       single-slide deck containing the panel

Style matches the existing presentation: dark-navy header bar, white
title text, light-blue page background, cyan accent stripe.
"""

from __future__ import annotations

import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Inches, Pt
except ImportError:
    sys.stderr.write("python-pptx is required: pip install python-pptx\n")
    sys.exit(1)

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))
import toy_data as td

PANEL_PATH = PROJECT_ROOT / "toy_datasets_panel.png"
PPTX_PATH = PROJECT_ROOT / "toy_datasets_intro.pptx"

# Match the existing deck's design tokens.
NAVY = RGBColor(0x1E, 0x27, 0x61)
CYAN = RGBColor(0x4F, 0xC3, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFF)
FONT = "Calibri"

# Hex equivalents for matplotlib.
NAVY_HEX = "#1E2761"
CYAN_HEX = "#4FC3F7"
LIGHT_BG_HEX = "#F0F4FF"


# ---------------------------------------------------------------------------
# Step 1 — render the panel of plots
# ---------------------------------------------------------------------------

def render_panel(out_path: Path) -> None:
    """2x3 grid: 5 dataset plots + 1 caption cell."""
    items = [
        ("exponential_multiplicative",
         td.make_exponential_multiplicative(),
         "y = a·exp(b·x)·η,  log η ~ N",
         "Best case: log-linearization is exact"),
        ("exponential_additive",
         td.make_exponential_additive(),
         "y = a·exp(b·x) + ε",
         "Failure mode: log distorts additive noise"),
        ("michaelis_menten",
         td.make_michaelis_menten(),
         "y = Vmax · x / (Km + x) + ε",
         "Saturating; tests reciprocal linearization"),
        ("logistic_growth",
         td.make_logistic_growth(),
         "y = L / (1 + exp(−k(x − x₀))) + ε",
         "S-curve; tests init sensitivity"),
        ("multivariable_nonlinear",
         td.make_multivariable_nonlinear(),
         "y = 2·exp(0.4·x₁) + 3·x₂^1.5 + 10/(1+x₃) + ε",
         "Designed for ZZU — no single linearization works"),
    ]

    fig, axes = plt.subplots(2, 3, figsize=(14, 7), constrained_layout=True)
    fig.patch.set_facecolor(LIGHT_BG_HEX)

    for ax, (name, bundle, formula, caption) in zip(axes.ravel()[:5], items):
        if bundle.X.shape[1] == 1:
            x = bundle.X.iloc[:, 0].values
            order = np.argsort(x)
            ax.scatter(x, bundle.y.values, s=14, alpha=0.55,
                       color="#888888", label="observed y")
            ax.plot(x[order], bundle.y_true.values[order],
                    color=NAVY_HEX, lw=2.2, label="true signal")
            ax.set_xlabel("x")
        else:
            # multivariable: show y vs x1 marginal (most informative single view)
            x = bundle.X["x1"].values
            ax.scatter(x, bundle.y.values, s=14, alpha=0.55,
                       color="#888888", label="observed y vs x₁")
            ax.set_xlabel("x₁  (also x₂, x₃ contribute)")
        ax.set_ylabel("y")
        ax.set_title(name, fontsize=11, color=NAVY_HEX, fontweight="bold")
        ax.set_facecolor("white")
        ax.grid(alpha=0.3)
        ax.legend(fontsize=8, loc="best")
        # caption below title
        ax.text(0.02, 0.98, formula, transform=ax.transAxes,
                fontsize=8, color="#444", va="top", family="monospace",
                bbox=dict(facecolor="white", alpha=0.85, edgecolor="none",
                          pad=2))
        ax.text(0.02, -0.22, caption, transform=ax.transAxes,
                fontsize=9, color=NAVY_HEX, va="top", style="italic")

    # 6th panel: text legend / "why these five"
    ax = axes[1, 2]
    ax.axis("off")
    ax.set_facecolor(LIGHT_BG_HEX)
    legend_text = (
        "Why these five?\n\n"
        "• Span the spectrum from\n"
        "  'log is exact' (favors\n"
        "  linearization) to 'no single\n"
        "  global transform' (favors ZZU)\n\n"
        "• Cover the three error\n"
        "  structures we care about:\n"
        "    – multiplicative lognormal\n"
        "    – additive Gaussian\n"
        "    – additive on saturating\n"
        "      / S-shaped responses\n\n"
        "• Multivariable case stress-\n"
        "  tests the hybrid screening\n"
        "  + warm-start workflow"
    )
    ax.text(0.02, 0.98, legend_text, transform=ax.transAxes,
            fontsize=10, color=NAVY_HEX, va="top", family="sans-serif")

    fig.suptitle("Synthetic Toy Datasets — Generators in toy_data.py",
                 fontsize=14, color=NAVY_HEX, fontweight="bold")
    fig.savefig(out_path, dpi=160, facecolor=LIGHT_BG_HEX)
    plt.close(fig)
    print(f"Wrote panel: {out_path}")


# ---------------------------------------------------------------------------
# Step 2 — wrap the panel in a one-slide .pptx
# ---------------------------------------------------------------------------

def _set_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(slide, x, y, w, h, color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if color is not None:
        _set_solid(sh, color)
    else:
        sh.fill.background()
        sh.line.fill.background()
    return sh


def _set_text(shape, text, *, size, bold=False, color=NAVY):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    run = p.add_run() if not p.runs else p.runs[0]
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def build_slide() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG

    # Header bar + title
    _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.65), color=NAVY)
    title_box = _add_rect(slide, Inches(0.4), Inches(0.1),
                          Inches(9.2), Inches(0.45))
    _set_text(title_box, "Synthetic Toy Datasets — What Each One Tests",
              size=22, bold=True, color=WHITE)

    # Cyan left accent stripe (matches deck)
    _add_rect(slide, Inches(0), Inches(0.65), Inches(0.07),
              Inches(5.625 - 0.65), color=CYAN)

    # Embed the panel image, centered horizontally below the header.
    from PIL import Image
    pix_w, pix_h = Image.open(PANEL_PATH).size
    aspect = pix_w / pix_h

    body_top = Inches(0.85)
    body_h = Inches(5.625 - 0.85 - 0.15)   # leave a thin bottom margin
    body_w = Inches(9.4)
    if body_w / aspect <= body_h:
        img_w = body_w
        img_h = Emu(int(body_w / aspect))
    else:
        img_h = body_h
        img_w = Emu(int(body_h * aspect))
    img_left = Inches(0.3) + Emu(int((Inches(9.4) - img_w) / 2))
    img_top = body_top + Emu(int((body_h - img_h) / 2))
    slide.shapes.add_picture(str(PANEL_PATH), img_left, img_top,
                             width=img_w, height=img_h)

    prs.save(PPTX_PATH)
    print(f"Wrote slide:  {PPTX_PATH}")


def main() -> None:
    render_panel(PANEL_PATH)
    build_slide()
    print("\nOpen toy_datasets_intro.pptx and copy the slide into your online deck.")


if __name__ == "__main__":
    main()
