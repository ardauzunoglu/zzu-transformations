"""
Append the nonlinear-methods, ZZU-workflow, benchmark, cost-analysis, and
selective-ZZU slides to ZZU_Transformations_Presentation.pptx.

Idempotent: if a slide with the same title already exists, this script
skips appending it. Re-running is therefore safe.

Usage:  python scripts/build_slides.py
"""

from __future__ import annotations

import sys
from pathlib import Path
from typing import Iterable, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Inches, Pt
except ImportError:
    sys.stderr.write(
        "python-pptx is required. Install with:\n"
        "    pip install python-pptx\n"
    )
    sys.exit(1)

from PIL import Image


# ---------------------------------------------------------------------------
# Paths and design tokens
# ---------------------------------------------------------------------------

PROJECT_ROOT = Path(__file__).resolve().parent.parent
PPTX_PATH = PROJECT_ROOT / "ZZU_Transformations_Presentation.pptx"
RESULTS_DIR = PROJECT_ROOT / "comparison_results"

# Color palette extracted from existing slides 1, 2, 6.
NAVY = RGBColor(0x1E, 0x27, 0x61)        # dark slide background, header bar, body text
NAVY_DEEP = RGBColor(0x21, 0x29, 0x5C)   # giant section-number watermark
CYAN = RGBColor(0x4F, 0xC3, 0xF7)        # accent stripe, "SECTION" tag
WHITE = RGBColor(0xFF, 0xFF, 0xFF)       # card fill, white title text
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFF)    # content-slide background
SUBTITLE = RGBColor(0xCA, 0xDC, 0xFC)    # cover subtitle
RED = RGBColor(0xE0, 0x52, 0x52)         # second-card accent stripe
FONT = "Calibri"

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _set_solid_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _set_no_fill(shape) -> None:
    shape.fill.background()
    shape.line.fill.background()


def _add_rect(slide, x, y, w, h, *, color: Optional[RGBColor] = None,
              shape_type=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    if color is None:
        _set_no_fill(sh)
    else:
        _set_solid_fill(sh, color)
    return sh


def _set_run(run, text: str, *, size: float, bold: bool = False,
             color: RGBColor = NAVY) -> None:
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_text(shape, text: str, *, size: float, bold: bool = False,
              color: RGBColor = NAVY) -> None:
    """Single paragraph, single run."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = None
    if not p.runs:
        run = p.add_run()
    else:
        run = p.runs[0]
    _set_run(run, text, size=size, bold=bold, color=color)


def _set_bullets(shape, lines: Iterable[str], *, size: float = 13,
                 color: RGBColor = NAVY, bullet: str = "•") -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = list(lines)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if not p.runs:
            run = p.add_run()
        else:
            run = p.runs[0]
        text = line if line.startswith(bullet) else f"{bullet}  {line}"
        _set_run(run, text, size=size, color=color)
        p.space_after = Pt(4)


def _set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _existing_titles(prs) -> List[str]:
    """Heuristic: first non-trivial text frame on each slide is the title."""
    titles = []
    for slide in prs.slides:
        title = ""
        for sh in slide.shapes:
            if sh.has_text_frame:
                txt = sh.text_frame.text.strip()
                if txt and not txt.isdigit() and txt.upper() != "SECTION":
                    title = txt.split("\n")[0].strip()
                    break
        titles.append(title)
    return titles


def _has_section_number(prs, number: str) -> bool:
    """Detect whether a section cover with the given giant number exists."""
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip() == number:
                return True
    return False


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_section_cover(prs, *, number: str, section_label: str,
                      title_lines: List[str], subtitle: str) -> None:
    """Replicate the section-cover layout from Slides 1 and 6."""
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, NAVY)

    # Cyan accent stripe on the left edge.
    _add_rect(slide, Inches(0), Inches(0), Inches(0.07), Inches(5.625),
              color=CYAN)

    # Giant section number on the right (acts as a darker watermark).
    num = _add_rect(slide, Inches(7.5), Inches(0.4), Inches(2.5), Inches(4.5))
    _set_text(num, number, size=180, bold=True, color=NAVY_DEEP)

    # "SECTION" tag.
    tag = _add_rect(slide, Inches(0.4), Inches(1.5), Inches(5.0), Inches(0.5))
    _set_text(tag, section_label, size=13, bold=True, color=CYAN)

    # Big title (one paragraph per line so the layout matches existing covers).
    title_box = _add_rect(slide, Inches(0.4), Inches(1.95),
                          Inches(7.0), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(title_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run() if not p.runs else p.runs[0]
        _set_run(run, line, size=54, bold=True, color=WHITE)

    # Subtitle.
    sub = _add_rect(slide, Inches(0.4), Inches(4.05),
                    Inches(7.5), Inches(0.9))
    tf = sub.text_frame
    tf.word_wrap = True
    for i, line in enumerate(subtitle.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run() if not p.runs else p.runs[0]
        _set_run(run, line, size=15, color=SUBTITLE)


def _add_content_header(slide, title: str) -> None:
    """Top header bar + white title text; matches Slide 2's pattern."""
    _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.65),
              color=NAVY)
    title_box = _add_rect(slide, Inches(0.4), Inches(0.1),
                          Inches(9.2), Inches(0.45))
    _set_text(title_box, title, size=22, bold=True, color=WHITE)


def _add_card(slide, x, y, w, h, *, accent_color: RGBColor,
              card_title: str, bullets: Iterable[str]) -> None:
    """White card with cyan/red accent stripe + heading + bullet list."""
    _add_rect(slide, x, y, w, h, color=WHITE)
    _add_rect(slide, x, y, Inches(0.07), h, color=accent_color)

    title_x = x + Inches(0.25)
    title_y = y + Inches(0.10)
    title_w = w - Inches(0.40)
    title = _add_rect(slide, title_x, title_y, title_w, Inches(0.4))
    _set_text(title, card_title, size=15, bold=True, color=NAVY)

    body_y = title_y + Inches(0.45)
    body_h = h - (body_y - y) - Inches(0.15)
    body = _add_rect(slide, title_x, body_y, title_w, body_h)
    _set_bullets(body, bullets, size=13, color=NAVY)


def add_one_card_slide(prs, *, title: str, card_title: str,
                       bullets: List[str],
                       accent_color: RGBColor = CYAN) -> None:
    """Content slide with a single full-width card."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_slide_bg(slide, LIGHT_BG)
    _add_content_header(slide, title)
    _add_card(slide, Inches(0.3), Inches(0.85),
              Inches(9.4), Inches(4.55),
              accent_color=accent_color,
              card_title=card_title, bullets=bullets)


def add_two_card_slide(prs, *, title: str,
                       left_card_title: str, left_bullets: List[str],
                       right_card_title: str, right_bullets: List[str],
                       left_accent: RGBColor = CYAN,
                       right_accent: RGBColor = RED) -> None:
    """Two side-by-side cards (matches Slide 2's layout)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_slide_bg(slide, LIGHT_BG)
    _add_content_header(slide, title)
    _add_card(slide, Inches(0.3), Inches(0.85),
              Inches(4.6), Inches(4.55),
              accent_color=left_accent,
              card_title=left_card_title, bullets=left_bullets)
    _add_card(slide, Inches(5.1), Inches(0.85),
              Inches(4.6), Inches(4.55),
              accent_color=right_accent,
              card_title=right_card_title, bullets=right_bullets)


def add_image_with_bullets_slide(
    prs, *, title: str, image_path: Path,
    bullets: List[str],
    image_position: str = "below",   # "below" or "left"
    bullets_card_title: str = "Highlights",
) -> None:
    """Image + a bullet card.  Layout depends on image aspect ratio."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_slide_bg(slide, LIGHT_BG)
    _add_content_header(slide, title)

    if not image_path.exists():
        # Fallback: bullet-only card.
        sys.stderr.write(f"warn: missing image {image_path}; bullets only\n")
        _add_card(slide, Inches(0.3), Inches(0.85),
                  Inches(9.4), Inches(4.55),
                  accent_color=CYAN,
                  card_title=bullets_card_title, bullets=bullets)
        return

    pix_w, pix_h = Image.open(image_path).size
    aspect = pix_w / pix_h

    body_top = Inches(0.85)
    body_h_total = Inches(4.55)

    if image_position == "left":
        # Image takes the left half; bullets on the right.
        max_w = Inches(4.6)
        max_h = body_h_total
        if max_w / aspect <= max_h:
            img_w = max_w
            img_h = Emu(int(max_w / aspect))
        else:
            img_h = max_h
            img_w = Emu(int(max_h * aspect))
        img_left = Inches(0.3)
        img_top = body_top + Emu(int((body_h_total - img_h) / 2))
        slide.shapes.add_picture(str(image_path), img_left, img_top,
                                 width=img_w, height=img_h)
        _add_card(slide, Inches(5.1), body_top,
                  Inches(4.6), body_h_total,
                  accent_color=CYAN,
                  card_title=bullets_card_title, bullets=bullets)
    else:
        # Image at top, full content width; bullet card below.
        max_w = Inches(9.4)
        # Limit image height to ~55% of body so bullets always have room.
        max_h_img = Inches(2.6)
        img_w = max_w
        img_h = Emu(int(max_w / aspect))
        if img_h > max_h_img:
            img_h = max_h_img
            img_w = Emu(int(max_h_img * aspect))
        img_left = Inches(0.3) + Emu(int((Inches(9.4) - img_w) / 2))
        img_top = body_top
        slide.shapes.add_picture(str(image_path), img_left, img_top,
                                 width=img_w, height=img_h)
        card_top = body_top + img_h + Inches(0.1)
        card_h = (Inches(0.85) + body_h_total) - card_top
        _add_card(slide, Inches(0.3), card_top,
                  Inches(9.4), card_h,
                  accent_color=CYAN,
                  card_title=bullets_card_title, bullets=bullets)


# ---------------------------------------------------------------------------
# Slide content (the actual deck additions)
# ---------------------------------------------------------------------------

def build_new_slides(prs) -> int:
    existing = _existing_titles(prs)

    def have(t: str) -> bool:
        # Exact first-paragraph match avoids false positives when an existing
        # slide's title merely contains the new title as a substring.
        return any(t == e for e in existing if e)

    added = 0

    # ---- Section 03: Nonlinear Optimization ----
    if not _has_section_number(prs, "03"):
        add_section_cover(
            prs,
            number="03",
            section_label="SECTION",
            title_lines=["Nonlinear", "Optimization"],
            subtitle=(
                "Direct minimization of SSE on the original scale —\n"
                "pure-NumPy gradient descent, Gauss-Newton, and BFGS"
            ),
        )
        added += 1

    if not have("Three Optimizers, One Interface"):
        add_two_card_slide(
            prs,
            title="Three Optimizers, One Interface",
            left_card_title="The Three Methods",
            left_bullets=[
                "Gradient Descent — vanilla GD with optional learning-rate decay; simple but slow",
                "Gauss-Newton — solve (JᵀJ + λI)δ = Jᵀr each step; fast on well-posed problems",
                "Self-activating Levenberg-Marquardt damping inside GN: λ=0 by default; rises on rejected steps",
                "BFGS — pure-NumPy inverse-Hessian update with backtracking Armijo line search; our default",
            ],
            right_card_title="Uniform API",
            right_bullets=[
                "reg.fit(X, y, theta_init) → self",
                "reg.predict(X_new) → ndarray (m,)",
                "reg.theta_, reg.converged_, reg.n_iter_, reg.fit_error_",
                "model_fn(X, theta) → (n,) is the only thing the user must supply",
                "Optional jacobian_fn; falls back to numerical Jacobian otherwise",
            ],
            left_accent=CYAN,
            right_accent=CYAN,
        )
        added += 1

    if not have("Numerical Jacobian + Convergence"):
        add_two_card_slide(
            prs,
            title="Numerical Jacobian + Convergence",
            left_card_title="Central Finite Differences",
            left_bullets=[
                "J[:, j] = (f(X, θ + h·eⱼ) − f(X, θ − h·eⱼ)) / (2h)",
                "Adaptive step: hⱼ = max(h, |θⱼ| · h) — robust near zero",
                "2p model evaluations per Jacobian (the cost bottleneck)",
                "Used automatically when jacobian_fn is not supplied",
            ],
            right_card_title="Convergence Criteria",
            right_bullets=[
                "GD: |SSEₖ − SSEₖ₊₁| / max(SSEₖ, ε) < tol",
                "GN: ‖δ‖ / (‖θ‖ + ε) < tol",
                "BFGS: ‖∇f(θ)‖ < tol  or  relative step size < tol",
                "Failures (singular matrix, NaN, overflow) caught and stored in reg.fit_error_",
                "Caller never sees a raised exception",
            ],
            left_accent=CYAN,
            right_accent=CYAN,
        )
        added += 1

    # ---- Section 04: ZZU Hybrid Workflow ----
    if not _has_section_number(prs, "04"):
        add_section_cover(
            prs,
            number="04",
            section_label="SECTION",
            title_lines=["The ZZU Hybrid", "Workflow"],
            subtitle="Screen → warm-start → bias correct",
        )
        added += 1

    if not have("ZZU in Three Steps"):
        add_one_card_slide(
            prs,
            title="ZZU in Three Steps",
            card_title="screen → warm-start → bias-correct",
            bullets=[
                "Step 1 — Screen: fit a suite of TransformedOLS models on a held-out validation split; rank by original-scale RMSE; refit the winner on full training data",
                "Step 2 — Warm start: invert the winning linearization to a nonlinear θ₀ via user-supplied coeff_to_init; hand θ₀ to BFGS / GN / GD",
                "Step 3 — Bias correct: add the mean training residual r̄ = ȳ − f(X, θ̂) as an additive correction at predict time",
                "Single fit() call exposes screening_table_, theta_init_used_, nonlinear_regressor_, train_residuals_ for diagnostics",
                "Failures gracefully captured in fit_error_ — pipeline never crashes",
            ],
            accent_color=CYAN,
        )
        added += 1

    if not have("The coeff_to_init Contract"):
        add_two_card_slide(
            prs,
            title="The coeff_to_init Contract",
            left_card_title="Why It's User-Supplied",
            left_bullets=[
                "Only the user knows how to map their chosen transform back to nonlinear parameters",
                "Example — y = a · exp(b · x):",
                "    log(y) = β₀ + β₁ · x  ⇒  a = exp(β₀),  b = β₁",
                "ZZU calls coeff_to_init(best_TransformedOLS_model)  →  θ₀",
                "Restrict the screening dict to transforms you can actually invert",
            ],
            right_card_title="Safety Nets",
            right_bullets=[
                "If coeff_to_init raises, ZZU falls back to fallback_theta_init",
                "If fallback is unset, ZZU uses np.ones(p_inferred)",
                "Failure is surfaced in fit_error_ but the pipeline still completes",
                "Mismatched coeff_to_init can drag the warm start *away* from the optimum — see warm-vs-cold analysis",
            ],
            left_accent=CYAN,
            right_accent=CYAN,
        )
        added += 1

    # ---- Synthetic benchmark ----
    if not have("Synthetic Benchmark — RMSE Across 5 Datasets"):
        add_image_with_bullets_slide(
            prs,
            title="Synthetic Benchmark — RMSE Across 5 Datasets",
            image_path=RESULTS_DIR / "rmse_by_method.png",
            image_position="left",
            bullets_card_title="Headline Numbers (10 splits)",
            bullets=[
                "exp_mult: log_smear wins (RMSE 7.58) — log linearization is exact under multiplicative noise",
                "exp_add: GD / GN / BFGS sweep top 3 (5.14) — log distorts additive noise",
                "Michaelis-Menten: BFGS / GN / ZZU tie (0.297, R² 0.97)",
                "logistic: ZZU edges nonlinear (3.065)",
                "multivariable: ZZU wins outright (4.65 vs 4.75) — the headline result",
                "ZZU wins or ties 3 of 5 datasets",
            ],
        )
        added += 1

    # ---- Cost analysis ----
    if not have("Cost vs. Accuracy"):
        add_image_with_bullets_slide(
            prs,
            title="Cost vs. Accuracy",
            image_path=RESULTS_DIR / "cost_pareto.png",
            image_position="below",
            bullets_card_title="Per-Family Cost (n = 120–500)",
            bullets=[
                "Linearized OLS: ~0.1 ms — only competitive on exp_mult",
                "Gauss-Newton / BFGS: ~1–10 ms — 10× slower, but lower RMSE on 4/5 datasets",
                "ZZU: ~3–30 ms — 3× overhead over BFGS for the screening phase",
                "Gradient descent: ~100 ms – 1 s; dominated everywhere",
                "ZZU appears on the Pareto frontier only on multivariable_nonlinear",
            ],
        )
        added += 1

    if not have("Warm-Start vs Cold-Start BFGS"):
        add_image_with_bullets_slide(
            prs,
            title="Warm-Start vs Cold-Start BFGS",
            image_path=RESULTS_DIR / "warm_vs_cold.png",
            image_position="below",
            bullets_card_title="Theory and Empirics Match",
            bullets=[
                "exp_multiplicative: warm cuts iterations 13 → 9 (~30%); log is the EXACT linearization",
                "exp_additive: warm INCREASES iterations 20 → 26; log distorts additive noise",
                "Concrete on exp_mult (true a=2, b=0.7): cold init b₀=0.10 (L²≈0.60); warm init b₀≈0.69 (L²≈0.01)",
                "Warm start helps only when the screened transform matches the noise structure",
            ],
        )
        added += 1

    # ---- Selective ZZU (future work) ----
    if not have("Future Work — Selective ZZU"):
        add_one_card_slide(
            prs,
            title="Future Work — Selective ZZU",
            card_title="Run ZZU only when the warm start will pay off",
            bullets=[
                "Motivation: warm-vs-cold analysis shows ZZU helps when the screened transform matches the noise, hurts when it doesn't",
                "Goal: gate ZZU on a confidence signal; otherwise fall back to cold-start nonlinear and skip the screening overhead",
                "Candidate signals — residual normality on the transformed scale (Shapiro-Wilk / Anderson-Darling)",
                "Candidate signals — screening RMSE ratio: best-transform val RMSE / identity val RMSE",
                "Candidate signals — log-likelihood gain over the identity baseline",
                "Status: in progress — being explored on the synthetic suite first",
            ],
            accent_color=CYAN,
        )
        added += 1

    return added


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    if not PPTX_PATH.exists():
        sys.stderr.write(f"error: {PPTX_PATH} not found\n")
        sys.exit(1)

    prs = Presentation(PPTX_PATH)
    n_before = len(prs.slides)
    print(f"Opened {PPTX_PATH.name} ({n_before} existing slides)")

    added = build_new_slides(prs)
    n_after = len(prs.slides)

    if added == 0:
        print("No new slides to add — all titles already present. Deck unchanged.")
        return

    prs.save(PPTX_PATH)
    print(f"Appended {added} slide(s); deck now has {n_after} slides.")
    print(f"Saved: {PPTX_PATH}")


if __name__ == "__main__":
    main()
